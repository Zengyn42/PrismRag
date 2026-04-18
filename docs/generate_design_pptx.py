"""Generate PrismRag design slide deck (pptx) — English edition.

High-level narration, dense visuals. Focus areas:
  - End-to-end indexing flow for a new note
  - End-to-end query flow
  - How LLMs consume the RAG / how the RAG is managed
  - Referenced prior art (GitHub projects, papers)
  - Integration with ZenithLoom

Output: PrismRag/docs/PrismRag-Design-Overview.pptx
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Inches, Pt

# ── Palette ─────────────────────────────────────────────────────────────
BG_DARK = RGBColor(0x0F, 0x1B, 0x2D)
BG_PANEL = RGBColor(0x1A, 0x2B, 0x45)
ACCENT = RGBColor(0x5E, 0xC5, 0xE8)  # cyan
ACCENT2 = RGBColor(0xF4, 0xB3, 0x5C)  # amber
ACCENT3 = RGBColor(0x8B, 0xE0, 0xA4)  # green
ACCENT4 = RGBColor(0xE8, 0x87, 0xC7)  # pink
TEXT_MAIN = RGBColor(0xE8, 0xEE, 0xF7)
TEXT_DIM = RGBColor(0x9A, 0xA8, 0xBE)
LINE = RGBColor(0x32, 0x4A, 0x6E)

# ── Layout ──────────────────────────────────────────────────────────────
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H

BLANK = prs.slide_layouts[6]


def add_bg(slide, color=BG_DARK):
    rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    rect.fill.solid()
    rect.fill.fore_color.rgb = color
    rect.line.fill.background()
    return rect


def add_text(
    slide, x, y, w, h, text,
    *, size=14, bold=False, color=TEXT_MAIN, align=PP_ALIGN.LEFT,
    font="Calibri",
):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    lines = text.split("\n") if isinstance(text, str) else text
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run()
        r.text = line
        r.font.name = font
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = color
    return tb


def add_panel(slide, x, y, w, h, *, fill=BG_PANEL, border=LINE):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    shape.adjustments[0] = 0.06
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = border
    shape.line.width = Pt(0.75)
    return shape


def add_title_bar(slide, title, subtitle=None):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(0.45), Inches(0.1), Inches(0.55))
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT
    bar.line.fill.background()

    add_text(slide, Inches(0.75), Inches(0.35), Inches(11.5), Inches(0.55),
             title, size=26, bold=True, color=TEXT_MAIN)
    if subtitle:
        add_text(slide, Inches(0.78), Inches(0.88), Inches(11.5), Inches(0.35),
                 subtitle, size=12, color=TEXT_DIM)

    rule = slide.shapes.add_connector(1, Inches(0.5), Inches(1.25), Inches(12.83), Inches(1.25))
    rule.line.color.rgb = LINE
    rule.line.width = Pt(0.75)


def add_footer(slide, idx, total):
    add_text(slide, Inches(0.5), Inches(7.05), Inches(6), Inches(0.3),
             "PrismRag · Design Overview · 2026-04-14", size=9, color=TEXT_DIM)
    add_text(slide, Inches(11.5), Inches(7.05), Inches(1.3), Inches(0.3),
             f"{idx} / {total}", size=9, color=TEXT_DIM, align=PP_ALIGN.RIGHT)


def arrow_right(slide, x, y, w, h, *, color=LINE):
    a = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, x, y, w, h)
    a.fill.solid()
    a.fill.fore_color.rgb = color
    a.line.fill.background()
    return a


def arrow_down(slide, x, y, w, h, *, color=LINE):
    a = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, x, y, w, h)
    a.fill.solid()
    a.fill.fore_color.rgb = color
    a.line.fill.background()
    return a


# ═══════════════════════════════════════════════════════════════════════
TOTAL = 17


def slide_01_cover():
    s = prs.slides.add_slide(BLANK)
    add_bg(s, BG_DARK)

    for i, (x, y, color) in enumerate([
        (Inches(10.5), Inches(1.5), ACCENT),
        (Inches(11.2), Inches(2.6), ACCENT2),
        (Inches(10.7), Inches(3.8), ACCENT3),
        (Inches(11.5), Inches(4.9), ACCENT4),
    ]):
        tri = s.shapes.add_shape(MSO_SHAPE.RIGHT_TRIANGLE, x, y, Inches(1.4), Inches(1.0))
        tri.fill.solid()
        tri.fill.fore_color.rgb = color
        tri.line.fill.background()
        tri.rotation = i * 15

    add_text(s, Inches(0.8), Inches(2.0), Inches(10), Inches(0.7),
             "PrismRag", size=54, bold=True, color=ACCENT)
    add_text(s, Inches(0.85), Inches(2.95), Inches(10), Inches(0.7),
             "A Graph-First RAG for Obsidian Vaults", size=26, bold=True, color=TEXT_MAIN)
    add_text(s, Inches(0.85), Inches(3.75), Inches(11), Inches(0.5),
             "v4.0 Design Overview — 5-pass index pipeline · zero vector search at query time · federated graphs",
             size=14, color=TEXT_DIM)

    add_text(s, Inches(0.85), Inches(5.7), Inches(10), Inches(0.4),
             "Foundation / ZenithLoom", size=12, color=ACCENT)
    add_text(s, Inches(0.85), Inches(6.05), Inches(10), Inches(0.4),
             "2026-04-14", size=11, color=TEXT_DIM)


def slide_02_what_is_it():
    s = prs.slides.add_slide(BLANK)
    add_bg(s)
    add_title_bar(s, "What is PrismRag",
                  "In one line: index the Obsidian vault as a graph, retrieve by graph traversal — no vector search at query time.")

    panels = [
        ("At Index Time", ACCENT, [
            "Read all markdown in the vault",
            "Extract wikilinks / tags / frontmatter",
            "Call Gemini Embedding to compute similarity",
            "Run Leiden community detection",
            "Persist one NetworkX graph + report",
        ]),
        ("At Query Time", ACCENT2, [
            "No vector search, no re-embedding",
            "Resolve entry nodes by label / alias / tag",
            "Traverse the graph with BFS / DFS",
            "Hard-stop on a token budget",
            "Return a coherent subgraph with provenance",
        ]),
        ("External Surface", ACCENT3, [
            "MCP server with 8 tools",
            "Any MCP-speaking agent can call it",
            "Multi-vault federation (ZenithLoom + EdenGateway + …)",
            "CLI: ingest / query / info / serve",
            "Human-readable HTML visualisation",
        ]),
    ]
    x0 = Inches(0.55)
    w = Inches(4.05)
    gap = Inches(0.1)
    for i, (title, color, bullets) in enumerate(panels):
        x = x0 + (w + gap) * i
        add_panel(s, x, Inches(1.55), w, Inches(5.3))
        hdr = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(1.55), w, Inches(0.55))
        hdr.adjustments[0] = 0.15
        hdr.fill.solid()
        hdr.fill.fore_color.rgb = color
        hdr.line.fill.background()
        add_text(s, x + Inches(0.25), Inches(1.65), w - Inches(0.5), Inches(0.4),
                 title, size=17, bold=True, color=BG_DARK)
        body = "\n".join(f"•  {b}" for b in bullets)
        add_text(s, x + Inches(0.3), Inches(2.4), w - Inches(0.5), Inches(4.3),
                 body, size=12, color=TEXT_MAIN)

    add_footer(s, 2, TOTAL)


def slide_03_design_dna():
    s = prs.slides.add_slide(BLANK)
    add_bg(s)
    add_title_bar(s, "Design DNA: School B",
                  "Embedding is used ONLY at index time to build edges. At query time the system touches zero vectors.")

    add_panel(s, Inches(0.55), Inches(1.55), Inches(12.25), Inches(2.0))
    add_text(s, Inches(0.75), Inches(1.7), Inches(12), Inches(0.4),
             "Three schools we evaluated", size=15, bold=True, color=ACCENT)

    schools = [
        ("A  Pure Graphify", "No embedding at all — LLM assigns every similarity edge",
         "expensive, slow for large vaults", ACCENT4),
        ("B  Graphify + Embedding (CHOSEN)",
         "Embedding builds similarity edges at index time; query is pure graph traversal",
         "v4.0", ACCENT3),
        ("C  Classical RAG", "Vector search + BM25 + RRF fusion at every query",
         "v3.2 legacy; abandoned", ACCENT2),
    ]
    for i, (name, detail, tag, color) in enumerate(schools):
        x = Inches(0.75) + Inches(4.05) * i
        add_text(s, x, Inches(2.1), Inches(3.9), Inches(0.4),
                 name, size=13, bold=True, color=color)
        add_text(s, x, Inches(2.5), Inches(3.9), Inches(0.7),
                 detail, size=11, color=TEXT_MAIN)
        add_text(s, x, Inches(3.15), Inches(3.9), Inches(0.3),
                 tag, size=10, color=TEXT_DIM)

    # Why B
    add_panel(s, Inches(0.55), Inches(3.75), Inches(6.0), Inches(3.25))
    add_text(s, Inches(0.75), Inches(3.9), Inches(5.6), Inches(0.4),
             "Why School B wins", size=15, bold=True, color=ACCENT)
    reasons = [
        "Traversal is explainable — every edge carries its source_pass",
        "Zero query-time API cost — no query embedding call",
        "Embedding participates in one write, one time",
        "Natural fit with Obsidian's [[wikilink]] culture",
        "Incrementally upgradable — Layer 2 LLM extraction",
        "  can be added later without breaking anything",
    ]
    add_text(s, Inches(0.75), Inches(4.3), Inches(5.6), Inches(2.7),
             "\n".join(f"•  {r}" for r in reasons), size=12, color=TEXT_MAIN)

    # Edge types
    add_panel(s, Inches(6.7), Inches(3.75), Inches(6.15), Inches(3.25))
    add_text(s, Inches(6.9), Inches(3.9), Inches(5.8), Inches(0.4),
             "Two kinds of edges", size=15, bold=True, color=ACCENT)
    add_text(s, Inches(6.9), Inches(4.35), Inches(5.8), Inches(0.35),
             "EXTRACTED  (deterministic, zero LLM)", size=12, bold=True, color=ACCENT3)
    add_text(s, Inches(6.9), Inches(4.7), Inches(5.8), Inches(0.8),
             "Parsed from wikilinks / tags / frontmatter\nsource_pass = ast",
             size=11, color=TEXT_MAIN)

    add_text(s, Inches(6.9), Inches(5.5), Inches(5.8), Inches(0.35),
             "INFERRED  (embedding similarity)", size=12, bold=True, color=ACCENT2)
    add_text(s, Inches(6.9), Inches(5.85), Inches(5.8), Inches(1.1),
             "Top-K cosine similarity + threshold\nDeduplicated against EXTRACTED pairs\nsource_pass = embedding",
             size=11, color=TEXT_MAIN)

    add_footer(s, 3, TOTAL)


def slide_04_prior_art():
    s = prs.slides.add_slide(BLANK)
    add_bg(s)
    add_title_bar(s, "Prior Art & Intellectual Lineage",
                  "PrismRag is a synthesis — not a from-scratch invention. Credit where due.")

    # Top row — direct inspiration
    add_panel(s, Inches(0.5), Inches(1.55), Inches(12.35), Inches(2.1))
    add_text(s, Inches(0.7), Inches(1.68), Inches(12), Inches(0.4),
             "Direct Inspiration", size=15, bold=True, color=ACCENT)
    add_text(s, Inches(0.7), Inches(2.1), Inches(12), Inches(0.4),
             "github.com/safishamsi/graphify",
             size=14, bold=True, color=ACCENT3, font="Consolas")
    add_text(s, Inches(0.7), Inches(2.5), Inches(12), Inches(1.1),
             "A radical proposal: once the graph is built, Leiden community detection runs on pure topology — no\n"
             "embedding space needed. Query is graph traversal, not vector search. PrismRag adopts this skeleton\n"
             "and relaxes it slightly to ‘School B’ (embedding only as an edge generator, never at query time).",
             size=11, color=TEXT_MAIN)

    # Bottom — three cards
    cards = [
        ("Leiden Algorithm", "Traag, Waltman & van Eck, 2019",
         "Improvement over Louvain — guarantees connected communities.\nNature Scientific Reports 9(1).",
         ACCENT2),
        ("Microsoft GraphRAG", "Contrast, not copy",
         "Uses LLM to extract entities + relations at index time. PrismRag rejects the LLM extraction step\n"
         "for cost and determinism.",
         ACCENT4),
        ("NetworkX + leidenalg", "Python ecosystem",
         "NetworkX for the graph model; python-igraph + leidenalg for partitioning. No custom graph engine.",
         ACCENT),
    ]
    x0 = Inches(0.5)
    w = Inches(4.05)
    gap = Inches(0.1)
    for i, (title, sub, body, color) in enumerate(cards):
        x = x0 + (w + gap) * i
        add_panel(s, x, Inches(3.85), w, Inches(3.1))
        bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, Inches(3.85), w, Inches(0.08))
        bar.fill.solid()
        bar.fill.fore_color.rgb = color
        bar.line.fill.background()
        add_text(s, x + Inches(0.2), Inches(4.05), w - Inches(0.4), Inches(0.4),
                 title, size=14, bold=True, color=color)
        add_text(s, x + Inches(0.2), Inches(4.45), w - Inches(0.4), Inches(0.35),
                 sub, size=10, color=TEXT_DIM)
        add_text(s, x + Inches(0.2), Inches(4.9), w - Inches(0.4), Inches(2.0),
                 body, size=11, color=TEXT_MAIN)

    add_footer(s, 4, TOTAL)


def slide_05_pipeline():
    s = prs.slides.add_slide(BLANK)
    add_bg(s)
    add_title_bar(s, "The 5-Pass Indexing Pipeline",
                  "One ingest run: vault → graph → report")

    stages = [
        ("Pass 1", "AST Extraction", "wikilinks · tags\nfrontmatter", "zero LLM\n100%", ACCENT3),
        ("Pass 2", "Media Extraction", "images · PDF · audio\n→ text nodes", "designed\nnot built", ACCENT4),
        ("Pass 3", "Embedding + Links", "Gemini Embedding 2\ncosine top-K", "edge engine\n100%", ACCENT3),
        ("Pass 4", "Leiden Clustering", "topology only\ngod-node ID", "100%", ACCENT3),
        ("Pass 5", "Persist + Report", "graph.json\nREPORT.md · HTML", "100%", ACCENT3),
    ]
    x0 = Inches(0.5)
    w = Inches(2.35)
    gap = Inches(0.12)
    y = Inches(1.75)
    for i, (tag, name, detail, status, color) in enumerate(stages):
        x = x0 + (w + gap) * i
        add_panel(s, x, y, w, Inches(3.2))
        strip = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, Inches(0.5))
        strip.adjustments[0] = 0.2
        strip.fill.solid()
        strip.fill.fore_color.rgb = color
        strip.line.fill.background()
        add_text(s, x, y + Inches(0.1), w, Inches(0.35), tag,
                 size=13, bold=True, color=BG_DARK, align=PP_ALIGN.CENTER)
        add_text(s, x + Inches(0.1), y + Inches(0.7), w - Inches(0.2), Inches(0.4),
                 name, size=14, bold=True, color=TEXT_MAIN, align=PP_ALIGN.CENTER)
        add_text(s, x + Inches(0.1), y + Inches(1.2), w - Inches(0.2), Inches(1.0),
                 detail, size=11, color=TEXT_DIM, align=PP_ALIGN.CENTER)
        add_text(s, x + Inches(0.1), y + Inches(2.5), w - Inches(0.2), Inches(0.55),
                 status, size=10, bold=True, color=color, align=PP_ALIGN.CENTER)
        if i < len(stages) - 1:
            arrow_right(s, x + w + Emu(10000), y + Inches(1.4),
                        gap - Emu(20000), Inches(0.3))

    add_panel(s, Inches(0.5), Inches(5.3), Inches(12.3), Inches(1.55))
    add_text(s, Inches(0.7), Inches(5.45), Inches(12), Inches(0.4),
             "Current completion: 4 of 5 passes production-ready (~80%)",
             size=15, bold=True, color=ACCENT)
    add_text(s, Inches(0.7), Inches(5.85), Inches(12), Inches(1.0),
             "Passes 1 / 3 / 4 / 5 run end-to-end. All 15 e2e tests pass. data/ already holds a 383 KB graph.json.\n"
             "Pass 2 (media) is fully designed but zero lines of code — the image/pdf/audio slots in NodeKind await a producer.",
             size=12, color=TEXT_MAIN)

    add_footer(s, 5, TOTAL)


def slide_06_ingest_flow():
    s = prs.slides.add_slide(BLANK)
    add_bg(s)
    add_title_bar(s, "End-to-End: Indexing a New Note",
                  "What happens the moment a new markdown file appears in the vault")

    # Left rail: vertical flow
    steps = [
        ("1  Detect",
         "vault_loader walks the vault. For each .md it reads mtime + SHA-256\nof the body. Compared against the previous ingest to classify:\nnew · changed · unchanged · deleted.",
         ACCENT),
        ("2  Parse frontmatter & body",
         "ast_extractor parses the YAML frontmatter (title, aliases, tags, type, …)\nand scans the body for [[wikilinks]], #tags, heading structure.\nStub nodes are created for link targets that do not yet exist.",
         ACCENT),
        ("3  Emit Pass-1 edges",
         "One EXTRACTED edge per wikilink / tag / category reference.\nsource_pass=ast, confidence=1.0. Deterministic — same input, same output.",
         ACCENT3),
        ("4  Embed (Pass 3a)",
         "embedder.py sends the (title + body) to Gemini Embedding 2.\nRate-limited at 0.5 s, body truncated at 30 k chars, 768-dim Matryoshka.\nVector is cached in LanceDB keyed by content_hash — no re-charge on re-ingest.",
         ACCENT2),
        ("5  Link similar (Pass 3b)",
         "similarity_linker walks all node pairs, keeps top-K above threshold,\nskips pairs already joined by an EXTRACTED edge. Emits\nsemantically_similar_to edges with source_pass=embedding.",
         ACCENT2),
        ("6  Re-cluster (Pass 4)",
         "Leiden partitions the updated graph. For incremental ingest we re-run\nLeiden globally — cheap at O(100s) of nodes. God-nodes (degree > threshold)\nare flagged; community density is recomputed.",
         ACCENT3),
        ("7  Persist & report (Pass 5)",
         "graph.json is written atomically, GRAPH_REPORT.md regenerated,\npyvis HTML refreshed. Incremental state table is updated for the next run.",
         ACCENT3),
    ]
    y = Inches(1.6)
    step_h = Inches(0.73)
    gap_h = Inches(0.05)
    for i, (tag, body, color) in enumerate(steps):
        # Step number badge
        badge = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.55), y + Inches(0.08), Inches(0.4), Inches(0.4))
        badge.fill.solid()
        badge.fill.fore_color.rgb = color
        badge.line.fill.background()
        add_text(s, Inches(0.55), y + Inches(0.13), Inches(0.4), Inches(0.3),
                 tag.split()[0], size=13, bold=True, color=BG_DARK, align=PP_ALIGN.CENTER)
        # Panel
        add_panel(s, Inches(1.1), y, Inches(11.75), step_h)
        add_text(s, Inches(1.3), y + Inches(0.06), Inches(11.3), Inches(0.3),
                 tag[3:].strip() if len(tag) > 3 else tag,
                 size=12, bold=True, color=color)
        add_text(s, Inches(1.3), y + Inches(0.32), Inches(11.3), Inches(0.45),
                 body, size=10, color=TEXT_MAIN)
        y += step_h + gap_h

    add_footer(s, 6, TOTAL)


def slide_07_data_model():
    s = prs.slides.add_slide(BLANK)
    add_bg(s)
    add_title_bar(s, "Data Model: Nodes and Edges",
                  "The K-atom triplet (name, attributes, relations) grounded in Python dataclasses")

    add_panel(s, Inches(0.5), Inches(1.55), Inches(6.1), Inches(5.25))
    add_text(s, Inches(0.7), Inches(1.7), Inches(5.8), Inches(0.4),
             "Node", size=17, bold=True, color=ACCENT)

    node_fields = [
        ("id / label", "Nm — identity and display label", ACCENT3),
        ("kind", "structural type:\nnote / tag / category / image / pdf / audio / section / block", TEXT_MAIN),
        ("source_file", "originating markdown path", TEXT_MAIN),
        ("content / content_hash", "body text + SHA-256 for change detection", TEXT_MAIN),
        ("tokens", "token count for budget trimming", TEXT_MAIN),
        ("frontmatter", "YAML metadata preserved verbatim", TEXT_MAIN),
        ("community_id", "Leiden cluster assignment", TEXT_MAIN),
        ("maturity · confidence · actionability",
         "Am — attributes filled by upstream Agent\n(seed/growing/mature · high/med/low · ref/decision/task)", ACCENT2),
    ]
    y = Inches(2.2)
    for name, desc, color in node_fields:
        add_text(s, Inches(0.75), y, Inches(5.7), Inches(0.28),
                 name, size=11, bold=True, color=color, font="Consolas")
        n_lines = desc.count("\n") + 1
        h = Inches(0.25 * n_lines + 0.08)
        add_text(s, Inches(0.9), y + Inches(0.25), Inches(5.55), h,
                 desc, size=10, color=TEXT_DIM)
        y += Inches(0.3) + h - Inches(0.05)

    add_panel(s, Inches(6.75), Inches(1.55), Inches(6.1), Inches(5.25))
    add_text(s, Inches(6.95), Inches(1.7), Inches(5.8), Inches(0.4),
             "Edge — Rm (relations)", size=17, bold=True, color=ACCENT)

    edge_fields = [
        ("source → target", "directed edge", TEXT_MAIN),
        ("relation", "links_to / tagged_with / in_category /\nsemantically_similar_to / …", TEXT_MAIN),
        ("confidence", "EXTRACTED / INFERRED / AMBIGUOUS", ACCENT3),
        ("confidence_score", "0.0 – 1.0 (similarity or curator score)", TEXT_MAIN),
        ("weight", "graph-algorithm weight", TEXT_MAIN),
        ("source_pass", "ast / media / embedding / llm\n(which pass produced this edge)", ACCENT2),
    ]
    y = Inches(2.2)
    for name, desc, color in edge_fields:
        add_text(s, Inches(7.0), y, Inches(5.7), Inches(0.28),
                 name, size=11, bold=True, color=color, font="Consolas")
        n_lines = desc.count("\n") + 1
        h = Inches(0.25 * n_lines + 0.08)
        add_text(s, Inches(7.15), y + Inches(0.25), Inches(5.55), h,
                 desc, size=10, color=TEXT_DIM)
        y += Inches(0.3) + h - Inches(0.05)

    add_text(s, Inches(6.95), Inches(6.25), Inches(5.8), Inches(0.5),
             "⚠  ontology_type (concept / entity / tool / …) is NOT yet implemented",
             size=11, bold=True, color=ACCENT4)

    add_footer(s, 7, TOTAL)


def slide_08_query_flow():
    s = prs.slides.add_slide(BLANK)
    add_bg(s)
    add_title_bar(s, "End-to-End: Answering a Query",
                  "Caller sends text; PrismRag returns a ranked, budget-bounded subgraph with provenance.")

    steps = [
        ("1  Receive request", "MCP tool call",
         "search_knowledge(query=\"...\", top_k=8, budget=4000)\narrives over the MCP protocol.\nOptional filters: namespace, kind, community_id."),
        ("2  Resolve entry nodes", "retrieve/entry.py",
         "Exact and substring match against node.label, node.id,\nfrontmatter.aliases, tag names. Multiple entry points\nare allowed; ambiguous matches return AMBIGUOUS flag."),
        ("3  Traverse", "bfs.py + dfs.py",
         "Expand neighbours bidirectionally. BFS for breadth-first\ntopic coverage; DFS for deep chains. Edge weight +\nsource_pass drive ranking — EXTRACTED edges preferred."),
        ("4  Token budget", "shared accumulator",
         "Running sum of Node.tokens. Stop when budget exhausted.\nFirst-node guarantee: the top entry is always included\neven if it alone would blow the budget."),
        ("5  Emit subgraph", "MCP response",
         "List of nodes (content + metadata + community_id) plus\nthe edges that connect them. Callers get a coherent,\nexplainable slice — not a bag of vector neighbours."),
    ]

    # Horizontal flow with 5 boxes
    x0 = Inches(0.5)
    w = Inches(2.4)
    gap = Inches(0.12)
    y = Inches(1.65)
    colors = [ACCENT, ACCENT2, ACCENT3, ACCENT4, ACCENT]
    for i, (tag, sub, body) in enumerate(steps):
        x = x0 + (w + gap) * i
        add_panel(s, x, y, w, Inches(3.6))
        add_text(s, x + Inches(0.2), y + Inches(0.15), w - Inches(0.4), Inches(0.4),
                 tag, size=13, bold=True, color=colors[i])
        add_text(s, x + Inches(0.2), y + Inches(0.55), w - Inches(0.4), Inches(0.3),
                 sub, size=10, color=TEXT_DIM, font="Consolas")
        add_text(s, x + Inches(0.2), y + Inches(0.95), w - Inches(0.4), Inches(2.5),
                 body, size=10, color=TEXT_MAIN)
        if i < 4:
            arrow_right(s, x + w + Emu(10000), y + Inches(1.6),
                        gap - Emu(20000), Inches(0.3))

    # Insight footer
    add_panel(s, Inches(0.5), Inches(5.5), Inches(12.35), Inches(1.4))
    add_text(s, Inches(0.7), Inches(5.65), Inches(12), Inches(0.4),
             "Why this matters",
             size=14, bold=True, color=ACCENT)
    add_text(s, Inches(0.7), Inches(6.05), Inches(12), Inches(0.85),
             "•  Zero API cost at query time — no query embedding, no re-rank model\n"
             "•  Every returned node is reachable from the entry via a named relation — traceable, not a vector-nearest black box\n"
             "•  Token budget is a hard ceiling — LLM context stays predictable regardless of graph size",
             size=11, color=TEXT_MAIN)

    add_footer(s, 8, TOTAL)


def slide_09_llm_usage():
    s = prs.slides.add_slide(BLANK)
    add_bg(s)
    add_title_bar(s, "How LLMs Use PrismRag",
                  "PrismRag is a tool, not a teammate. LLMs (Hani, Jei, Asa, Apex Coder) call it via MCP.")

    # Left: call pattern
    add_panel(s, Inches(0.5), Inches(1.55), Inches(6.1), Inches(5.3))
    add_text(s, Inches(0.7), Inches(1.7), Inches(5.8), Inches(0.4),
             "Typical agent interaction", size=15, bold=True, color=ACCENT)
    pattern = [
        ("User →", "\"What did we decide about fresh_per_call session mode?\""),
        ("Agent", "decides this is a vault question →"),
        ("→ tool call", "search_knowledge(query=\"fresh_per_call\", budget=3000)"),
        ("PrismRag", "entry resolve → BFS traverse → subgraph"),
        ("→ response", "~8 nodes with content, edges, community labels"),
        ("Agent", "synthesises an answer citing source_file paths"),
        ("User ←", "answer with citations to vault notes"),
    ]
    y = Inches(2.2)
    for lbl, txt in pattern:
        add_text(s, Inches(0.75), y, Inches(1.3), Inches(0.3),
                 lbl, size=11, bold=True, color=ACCENT3, font="Consolas")
        add_text(s, Inches(2.1), y, Inches(4.4), Inches(0.3),
                 txt, size=11, color=TEXT_MAIN, font="Consolas")
        y += Inches(0.35)

    add_text(s, Inches(0.75), Inches(5.35), Inches(5.6), Inches(1.4),
             "Key idea: the agent stays in charge of synthesis and natural-language output.\n"
             "PrismRag only supplies the right subgraph — no prose generation, no opinion.",
             size=11, color=TEXT_DIM)

    # Right: tool menu
    add_panel(s, Inches(6.75), Inches(1.55), Inches(6.1), Inches(5.3))
    add_text(s, Inches(6.95), Inches(1.7), Inches(5.8), Inches(0.4),
             "The 8 MCP tools an agent sees", size=15, bold=True, color=ACCENT)

    tools = [
        ("search_knowledge", "entry resolve + BFS → relevant subgraph"),
        ("explain_node", "a single node + its neighbours + its community"),
        ("trace_path", "shortest path / relation chain between two nodes"),
        ("list_communities", "catalog of Leiden clusters + representatives"),
        ("explore_community", "drill into one cluster: members, density, god-node"),
        ("list_namespaces", "in federation: available vault namespaces"),
        ("read_note", "raw markdown of a specific path"),
        ("write_note", "write back to vault (CAS + audit — WIP)"),
    ]
    y = Inches(2.2)
    for name, desc in tools:
        add_text(s, Inches(7.0), y, Inches(3.2), Inches(0.3),
                 name, size=11, bold=True, color=ACCENT3, font="Consolas")
        add_text(s, Inches(10.2), y, Inches(2.65), Inches(0.3),
                 desc, size=10, color=TEXT_MAIN)
        y += Inches(0.38)

    add_text(s, Inches(6.95), Inches(5.6), Inches(5.8), Inches(1.2),
             "Agents compose these tools. For example, Jei (Knowledge Curator)\n"
             "calls search_knowledge → explain_node → write_note in sequence\n"
             "when promoting a dialogue insight into a durable note.",
             size=11, color=TEXT_DIM)

    add_footer(s, 9, TOTAL)


def slide_10_internal_llm():
    s = prs.slides.add_slide(BLANK)
    add_bg(s)
    add_title_bar(s, "Is There an LLM Inside PrismRag?",
                  "Short answer: no. Embedding yes, generative LLM no. By deliberate design.")

    # Left: today
    add_panel(s, Inches(0.5), Inches(1.55), Inches(6.1), Inches(5.3))
    add_text(s, Inches(0.7), Inches(1.7), Inches(5.8), Inches(0.4),
             "Today — the components PrismRag calls", size=15, bold=True, color=ACCENT)
    items = [
        ("Gemini Embedding 2 API", "index time only — per-note vector",
         "one API call per ingest per node; cached in LanceDB by hash", ACCENT2),
        ("leidenalg", "clustering", "no API, pure CPU algorithm", ACCENT3),
        ("NetworkX", "graph model + BFS / DFS",
         "no API, in-memory", ACCENT3),
        ("pyvis", "HTML visualisation", "no API, static output", ACCENT3),
        ("No generative LLM", "—",
         "Pass 1 is regex-deterministic; Pass 3 is embedding only.\n"
         "Nothing reads a note and calls Claude or Gemini-Pro.", ACCENT4),
    ]
    y = Inches(2.2)
    for name, kind, body, color in items:
        add_text(s, Inches(0.7), y, Inches(5.8), Inches(0.3),
                 name, size=12, bold=True, color=color)
        add_text(s, Inches(0.7), y + Inches(0.28), Inches(5.8), Inches(0.3),
                 kind, size=10, color=TEXT_DIM, font="Consolas")
        n = body.count("\n") + 1
        h = Inches(0.22 * n + 0.15)
        add_text(s, Inches(0.7), y + Inches(0.56), Inches(5.8), h,
                 body, size=10, color=TEXT_MAIN)
        y += Inches(0.62) + h - Inches(0.15)

    # Right: design rationale + future
    add_panel(s, Inches(6.75), Inches(1.55), Inches(6.1), Inches(5.3))
    add_text(s, Inches(6.95), Inches(1.7), Inches(5.8), Inches(0.4),
             "Why no LLM in the loop", size=15, bold=True, color=ACCENT)
    reasons = [
        ("Determinism", "Same vault, same graph — always. LLM extraction would make every ingest non-reproducible."),
        ("Cost", "An LLM per note per ingest scales badly. Embedding is one API call with a 24-hour cache."),
        ("Explainability", "Every edge traces to a literal wikilink or a cosine number. No ‘why did the LLM say this?’"),
        ("Separation of concerns", "LLMs live outside PrismRag and call it. Putting an LLM inside would blur who is deciding what."),
    ]
    y = Inches(2.2)
    for t, b in reasons:
        add_text(s, Inches(6.95), y, Inches(5.8), Inches(0.3),
                 t, size=12, bold=True, color=ACCENT3)
        add_text(s, Inches(6.95), y + Inches(0.28), Inches(5.8), Inches(0.7),
                 b, size=10, color=TEXT_MAIN)
        y += Inches(0.85)

    add_text(s, Inches(6.95), Inches(5.85), Inches(5.8), Inches(0.9),
             "Future (feature-flagged, default OFF): Layer 2 LLM extraction\n"
             "for ontology_type classification on ambiguous nodes. Will be\n"
             "opt-in and will not change the core invariant.",
             size=10, color=TEXT_DIM)

    add_footer(s, 10, TOTAL)


def slide_11_management():
    s = prs.slides.add_slide(BLANK)
    add_bg(s)
    add_title_bar(s, "How the RAG is Managed",
                  "Operational surface: commands, triggers, and the lifecycle of a graph")

    # Top row — three mgmt surfaces
    surfaces = [
        ("CLI", ACCENT, [
            "prism-rag ingest <vault>",
            "  Full or incremental (auto-detected)",
            "prism-rag query \"…\"",
            "  Dry-run the MCP search path",
            "prism-rag info",
            "  Stats · last ingest · namespaces",
            "prism-rag serve",
            "  Start MCP server (WIP)",
        ]),
        ("Config", ACCENT2, [
            "config.py → PrismRagSettings",
            "PRISM_VAULT_PATHS, PRISM_DATA_DIR",
            "PRISM_GEMINI_API_KEY",
            "PRISM_SIM_THRESHOLD (default 0.72)",
            "PRISM_TOP_K_SIMILARITY (default 5)",
            "PRISM_DEFAULT_BUDGET (tokens)",
            "GraphSource list for federation",
        ]),
        ("Automation", ACCENT3, [
            "systemd timer: nightly ingest",
            "Obsidian file watcher (future)",
            "incremental.py reuses the LanceDB",
            "  embedding cache — only new/changed",
            "  notes pay the embedding cost",
            "Idempotent: same vault state",
            "  → same graph.json byte-for-byte",
        ]),
    ]
    x0 = Inches(0.5)
    w = Inches(4.1)
    gap = Inches(0.1)
    for i, (title, color, items) in enumerate(surfaces):
        x = x0 + (w + gap) * i
        add_panel(s, x, Inches(1.55), w, Inches(3.5))
        hdr = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(1.55), w, Inches(0.5))
        hdr.adjustments[0] = 0.2
        hdr.fill.solid()
        hdr.fill.fore_color.rgb = color
        hdr.line.fill.background()
        add_text(s, x + Inches(0.2), Inches(1.62), w - Inches(0.4), Inches(0.4),
                 title, size=14, bold=True, color=BG_DARK)
        add_text(s, x + Inches(0.25), Inches(2.2), w - Inches(0.45), Inches(2.8),
                 "\n".join(items), size=10, color=TEXT_MAIN, font="Consolas")

    # Bottom — lifecycle
    add_panel(s, Inches(0.5), Inches(5.2), Inches(12.35), Inches(1.75))
    add_text(s, Inches(0.7), Inches(5.35), Inches(12), Inches(0.4),
             "Graph lifecycle", size=14, bold=True, color=ACCENT)

    life = [
        ("bootstrap", "first ingest — full vault scan, all nodes embedded"),
        ("incremental", "diff mtime + hash — only changed notes re-embedded"),
        ("re-cluster", "Leiden re-runs globally each ingest (cheap at O(100s))"),
        ("serve", "load graph.json into memory, answer MCP calls"),
        ("evolve", "model upgrade → versioned re-embed flow (planned)"),
    ]
    x = Inches(0.7)
    for lbl, desc in life:
        badge = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(5.85), Inches(1.4), Inches(0.35))
        badge.adjustments[0] = 0.4
        badge.fill.solid()
        badge.fill.fore_color.rgb = ACCENT
        badge.line.fill.background()
        add_text(s, x, Inches(5.9), Inches(1.4), Inches(0.3),
                 lbl, size=10, bold=True, color=BG_DARK, align=PP_ALIGN.CENTER)
        add_text(s, x, Inches(6.28), Inches(2.3), Inches(0.55),
                 desc, size=9, color=TEXT_MAIN)
        x += Inches(2.45)

    add_footer(s, 11, TOTAL)


def slide_12_federation():
    s = prs.slides.add_slide(BLANK)
    add_bg(s)
    add_title_bar(s, "Federated Graphs: Many Vaults, One View",
                  "Each vault ingests independently; a runtime layer stitches them into a unified queryable graph.")

    add_panel(s, Inches(0.5), Inches(1.55), Inches(5.5), Inches(5.3))
    add_text(s, Inches(0.7), Inches(1.7), Inches(5.2), Inches(0.4),
             "How it works", size=15, bold=True, color=ACCENT)
    bullets = [
        "Each vault ingests on its own schedule\nand produces its own graph.json",
        "FederatedGraph loads them at runtime\nand namespaces every node ID",
        "Cross-vault wikilinks and common labels\nbecome bridge edges",
        "Queries can filter by namespace\nor span all vaults",
        "Updating one vault never forces re-ingest\nof another — indexing is independent",
    ]
    add_text(s, Inches(0.7), Inches(2.15), Inches(5.2), Inches(4.6),
             "\n\n".join(f"•  {b}" for b in bullets),
             size=12, color=TEXT_MAIN)

    add_panel(s, Inches(6.2), Inches(1.55), Inches(6.65), Inches(5.3))
    add_text(s, Inches(6.4), Inches(1.7), Inches(6.3), Inches(0.4),
             "Example deployment", size=15, bold=True, color=ACCENT)

    vaults = [
        ("ZenithLoom", "designs · decisions · architecture", Inches(6.5), Inches(2.3), ACCENT),
        ("EdenGateway", "runtime data · session logs", Inches(6.5), Inches(3.3), ACCENT2),
        ("Resource", "external papers · references", Inches(6.5), Inches(4.3), ACCENT3),
    ]
    for name, desc, x, y, color in vaults:
        box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(3.1), Inches(0.75))
        box.adjustments[0] = 0.15
        box.fill.solid()
        box.fill.fore_color.rgb = BG_DARK
        box.line.color.rgb = color
        box.line.width = Pt(1.5)
        add_text(s, x + Inches(0.2), y + Inches(0.1), Inches(2.8), Inches(0.3),
                 name, size=13, bold=True, color=color)
        add_text(s, x + Inches(0.2), y + Inches(0.4), Inches(2.8), Inches(0.3),
                 desc, size=10, color=TEXT_DIM)

    fed = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                             Inches(10.0), Inches(3.3), Inches(2.7), Inches(0.75))
    fed.adjustments[0] = 0.15
    fed.fill.solid()
    fed.fill.fore_color.rgb = ACCENT4
    fed.line.fill.background()
    add_text(s, Inches(10.1), Inches(3.4), Inches(2.5), Inches(0.3),
             "FederatedGraph", size=13, bold=True, color=BG_DARK)
    add_text(s, Inches(10.1), Inches(3.7), Inches(2.5), Inches(0.3),
             "unified query view", size=10, color=BG_DARK)

    for y in [Inches(2.65), Inches(3.65), Inches(4.65)]:
        arr = s.shapes.add_connector(1, Inches(9.65), y, Inches(10.0), Inches(3.7))
        arr.line.color.rgb = LINE
        arr.line.width = Pt(1.2)

    add_text(s, Inches(6.4), Inches(5.4), Inches(6.3), Inches(1.3),
             "Implementation: store/federated.py (265 LOC)\n"
             "Tests: test_federated.py (28 cases) + multi-graph scenarios in e2e",
             size=11, color=TEXT_DIM, font="Consolas")

    add_footer(s, 12, TOTAL)


def slide_13_zenithloom():
    s = prs.slides.add_slide(BLANK)
    add_bg(s)
    add_title_bar(s, "Integration with ZenithLoom",
                  "PrismRag is ZenithLoom's long-term knowledge substrate — but the two evolve on different timelines")

    # Left: relationship
    add_panel(s, Inches(0.5), Inches(1.55), Inches(6.1), Inches(5.3))
    add_text(s, Inches(0.7), Inches(1.7), Inches(5.8), Inches(0.4),
             "Relationship", size=15, bold=True, color=ACCENT)
    rel = [
        ("ZenithLoom", "multi-agent framework — Hani, Jei, Asa, Apex Coder",
         "LangGraph-based, one LLM per role, runtime for everything", ACCENT),
        ("Vault", "the shared memory surface — Obsidian markdown on disk",
         "owned by Jei (Knowledge Curator); Hani et al. read via MCP", ACCENT3),
        ("PrismRag", "the indexer + query engine over that vault",
         "offline tool today; MCP-callable once `serve` ships", ACCENT2),
        ("How they meet", "via MCP",
         "An agent calls a PrismRag tool the same way it calls Obsidian MCP\n"
         "or any other tool — no special wiring, no coupling.", ACCENT4),
    ]
    y = Inches(2.2)
    for t, k, b, color in rel:
        add_text(s, Inches(0.7), y, Inches(5.8), Inches(0.3),
                 t, size=12, bold=True, color=color)
        add_text(s, Inches(0.7), y + Inches(0.28), Inches(5.8), Inches(0.3),
                 k, size=10, color=TEXT_DIM)
        n = b.count("\n") + 1
        h = Inches(0.22 * n + 0.1)
        add_text(s, Inches(0.7), y + Inches(0.56), Inches(5.8), h,
                 b, size=10, color=TEXT_MAIN)
        y += Inches(0.6) + h - Inches(0.1)

    # Right: phase alignment
    add_panel(s, Inches(6.75), Inches(1.55), Inches(6.1), Inches(5.3))
    add_text(s, Inches(6.95), Inches(1.7), Inches(5.8), Inches(0.4),
             "Phase alignment — two different clocks", size=15, bold=True, color=ACCENT)
    add_text(s, Inches(6.95), Inches(2.1), Inches(5.8), Inches(0.4),
             "ZenithLoom Vault phase (data model evolution)",
             size=11, bold=True, color=ACCENT3)
    zl_phases = [
        ("Phase 1 (now)", "raw markdown + wikilinks, no knowledge_id"),
        ("Phase 2", "add knowledge_id, relations:, lifecycle status"),
        ("Phase 3", "embedding for retrieval only"),
        ("Phase 4", "automatic edge discovery (≈ where PrismRag lives today)"),
    ]
    y = Inches(2.5)
    for p, d in zl_phases:
        add_text(s, Inches(7.0), y, Inches(1.85), Inches(0.3),
                 p, size=10, bold=True, color=ACCENT, font="Consolas")
        add_text(s, Inches(8.9), y, Inches(3.95), Inches(0.3),
                 d, size=10, color=TEXT_MAIN)
        y += Inches(0.32)

    add_text(s, Inches(6.95), Inches(4.0), Inches(5.8), Inches(0.4),
             "PrismRag phase (tool-build evolution)",
             size=11, bold=True, color=ACCENT2)
    pr_phases = [
        ("Phase 1 MVP", "AST + Leiden + persistence (done)"),
        ("Phase 2", "embedding + similarity + MCP + federation (~85% done)"),
        ("Phase 3", "media extraction (Pass 2) — not started"),
    ]
    y = Inches(4.4)
    for p, d in pr_phases:
        add_text(s, Inches(7.0), y, Inches(1.85), Inches(0.3),
                 p, size=10, bold=True, color=ACCENT, font="Consolas")
        add_text(s, Inches(8.9), y, Inches(3.95), Inches(0.3),
                 d, size=10, color=TEXT_MAIN)
        y += Inches(0.32)

    add_text(s, Inches(6.95), Inches(5.55), Inches(5.8), Inches(1.3),
             "Key tension: PrismRag indexes at file granularity today.\n"
             "Once Vault reaches Phase 2 (knowledge_id per atomic note),\n"
             "the node unit should shift from ‘file’ to ‘knowledge_id’.\n"
             "This is the single biggest integration change coming.",
             size=10, color=TEXT_DIM)

    add_footer(s, 13, TOTAL)


def slide_14_module_map():
    s = prs.slides.add_slide(BLANK)
    add_bg(s)
    add_title_bar(s, "Module Map & LOC",
                  "~3.5 k lines of main code + ~1.4 k lines of tests")

    modules = [
        ("ingest/", "AST · embedding · similarity · incremental", "~400", "🟢  80%"),
        ("store/", "KnowledgeGraph · embedding cache · federation", "~540", "🟢  90%"),
        ("retrieve/", "entry resolve · BFS · DFS", "~460", "🟢  90%"),
        ("cluster/", "Leiden · god-node · density", "154", "🟢  100%"),
        ("report/", "GRAPH_REPORT.md · pyvis HTML", "~340", "🟢  100%"),
        ("vault_ops/", "CAS · path sandbox · audit log", "~480", "🟡  60%"),
        ("mcp_server/", "8 MCP tools + transport", "637", "🟡  85%"),
        ("config.py", "GraphSource · federation settings", "107", "🟢  100%"),
        ("cli.py", "ingest / query / info / serve", "251", "🟡  90%"),
    ]

    x0 = Inches(0.5)
    w_total = Inches(12.35)
    col_w = [Inches(2.2), Inches(6.0), Inches(1.6), Inches(2.55)]
    y = Inches(1.65)

    hdr = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x0, y, w_total, Inches(0.45))
    hdr.fill.solid()
    hdr.fill.fore_color.rgb = BG_PANEL
    hdr.line.fill.background()
    headers = ["Module", "Responsibility", "LOC", "Status"]
    x = x0
    for i, htxt in enumerate(headers):
        add_text(s, x + Inches(0.15), y + Inches(0.1), col_w[i] - Inches(0.2), Inches(0.3),
                 htxt, size=12, bold=True, color=ACCENT)
        x += col_w[i]
    y += Inches(0.5)

    for i, (name, desc, loc, status) in enumerate(modules):
        if i % 2 == 0:
            row = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x0, y, w_total, Inches(0.45))
            row.fill.solid()
            row.fill.fore_color.rgb = BG_PANEL
            row.line.fill.background()
        x = x0
        for j, (val, col) in enumerate([
            (name, ACCENT3),
            (desc, TEXT_MAIN),
            (loc, TEXT_DIM),
            (status, TEXT_MAIN),
        ]):
            font = "Consolas" if j in (0, 2) else "Calibri"
            add_text(s, x + Inches(0.15), y + Inches(0.1), col_w[j] - Inches(0.2), Inches(0.3),
                     val, size=11, bold=(j == 0), color=col, font=font)
            x += col_w[j]
        y += Inches(0.45)

    add_text(s, Inches(0.5), Inches(6.4), Inches(12.3), Inches(0.5),
             "59 tests in 4 files: phase1_mvp (16) · embedding_store (4) · federated (28) · e2e_full (15)",
             size=11, color=TEXT_DIM, font="Consolas")

    add_footer(s, 14, TOTAL)


def slide_15_kg_alignment():
    s = prs.slides.add_slide(BLANK)
    add_bg(s)
    add_title_bar(s, "K-atom Alignment",
                  "How far the knowledge-element triplet (Nm, Am, Rm) is realised today")

    rows = [
        ("Nm — name", "Node.id + Node.label", "✅ complete", ACCENT3),
        ("Rm — relations", "Edge + confidence + source_pass", "✅ complete", ACCENT3),
        ("Am — attributes", "maturity / confidence / actionability",
         "🟨 schema in place, filled by upstream Agent", ACCENT2),
        ("ontology_type", "concept / entity / tool / process / …",
         "❌ missing — only structural kinds today", ACCENT4),
        ("Layer-1 extraction", "wikilinks + tags + frontmatter",
         "✅ complete (deterministic)", ACCENT3),
        ("Layer-2 LLM extraction", "Claude Haiku ontology classification",
         "❌ not implemented (and not an early priority)", ACCENT4),
        ("Alias merging", "frontmatter.aliases → one node",
         "🟨 aliases are read; actual merge behaviour needs verification", ACCENT2),
        ("Feature flags", "[features] to gate extensions",
         "❌ not yet built", ACCENT4),
    ]
    x0 = Inches(0.5)
    w_total = Inches(12.35)
    col_w = [Inches(2.4), Inches(4.5), Inches(5.45)]
    y = Inches(1.65)

    hdr = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x0, y, w_total, Inches(0.45))
    hdr.fill.solid()
    hdr.fill.fore_color.rgb = BG_PANEL
    hdr.line.fill.background()
    headers = ["K-atom dimension", "PrismRag mapping", "Status"]
    x = x0
    for i, htxt in enumerate(headers):
        add_text(s, x + Inches(0.15), y + Inches(0.1), col_w[i] - Inches(0.2), Inches(0.3),
                 htxt, size=12, bold=True, color=ACCENT)
        x += col_w[i]
    y += Inches(0.5)

    for i, (dim, mapping, status, color) in enumerate(rows):
        if i % 2 == 0:
            row = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x0, y, w_total, Inches(0.5))
            row.fill.solid()
            row.fill.fore_color.rgb = BG_PANEL
            row.line.fill.background()
        add_text(s, x0 + Inches(0.15), y + Inches(0.12), col_w[0] - Inches(0.2), Inches(0.3),
                 dim, size=11, bold=True, color=TEXT_MAIN)
        add_text(s, x0 + col_w[0] + Inches(0.15), y + Inches(0.12),
                 col_w[1] - Inches(0.2), Inches(0.3),
                 mapping, size=10, color=TEXT_DIM, font="Consolas")
        add_text(s, x0 + col_w[0] + col_w[1] + Inches(0.15), y + Inches(0.12),
                 col_w[2] - Inches(0.2), Inches(0.3),
                 status, size=10, bold=True, color=color)
        y += Inches(0.5)

    add_footer(s, 15, TOTAL)


def slide_16_roadmap():
    s = prs.slides.add_slide(BLANK)
    add_bg(s)
    add_title_bar(s, "Roadmap",
                  "Ranked P1 → P4. P1 is the blocker for PrismRag being useful to live agents.")

    groups = [
        ("P1  Blocking", ACCENT4, [
            "`prism-rag serve` entry point (stdio / SSE transport)",
            "Pass 2 media extraction (minimum: PDF)",
        ]),
        ("P2  Correctness", ACCENT2, [
            "tiktoken precise token count (replace len() // 4)",
            "Alias → node merge at ingest",
            "ontology_type field + frontmatter mapping",
            "vault_ops write path: CAS + audit + atomicity",
        ]),
        ("P3  Alignment & Extension", ACCENT, [
            "Default Am fill policy (frontmatter → maturity)",
            "Adopt ZenithLoom knowledge_id granularity (when Vault Phase 2 lands)",
            "Activate source_pass=media / llm edge types",
        ]),
        ("P4  Quality & Observability", ACCENT3, [
            "CLI integration tests",
            "Dangling wikilinks report",
            "Feature-flag infrastructure",
            "Embedding-model migration / re-embed flow",
        ]),
    ]

    x0 = Inches(0.5)
    w = Inches(6.1)
    h = Inches(2.6)
    gap_x = Inches(0.15)
    gap_y = Inches(0.2)
    for i, (title, color, items) in enumerate(groups):
        col = i % 2
        row = i // 2
        x = x0 + (w + gap_x) * col
        y = Inches(1.65) + (h + gap_y) * row
        add_panel(s, x, y, w, h)
        hdr = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(1.9), Inches(0.5))
        hdr.adjustments[0] = 0.25
        hdr.fill.solid()
        hdr.fill.fore_color.rgb = color
        hdr.line.fill.background()
        add_text(s, x, y + Inches(0.1), Inches(1.9), Inches(0.32),
                 title, size=13, bold=True, color=BG_DARK, align=PP_ALIGN.CENTER)
        body = "\n".join(f"•  {it}" for it in items)
        add_text(s, x + Inches(0.25), y + Inches(0.65), w - Inches(0.5), h - Inches(0.8),
                 body, size=12, color=TEXT_MAIN)

    add_footer(s, 16, TOTAL)


def slide_17_summary():
    s = prs.slides.add_slide(BLANK)
    add_bg(s)
    add_title_bar(s, "Summary",
                  "Where PrismRag stands and what it asks of you")

    insights = [
        ("What works today",
         "4 of 5 passes are production-ready. The graph builds, queries run, reports render, e2e tests pass.",
         ACCENT3),
        ("The core bet — School B",
         "Embedding only at index time; query-time traversal is pure, explainable, and free of API calls.",
         ACCENT),
        ("The next milestone",
         "Ship `serve` + a minimum Pass 2 (PDF) and PrismRag becomes a live tool for every ZenithLoom agent.",
         ACCENT2),
        ("Coupling with ZenithLoom",
         "PrismRag indexes files today. When Vault enters Phase 2 (knowledge_id), switch the node granularity.",
         ACCENT4),
        ("What NOT to add",
         "Layer 2 LLM extraction and premature ontology inference — they violate the simplicity of School B.",
         ACCENT4),
    ]
    y = Inches(1.65)
    for title, body, color in insights:
        add_panel(s, Inches(0.5), y, Inches(12.35), Inches(0.98))
        bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), y, Inches(0.12), Inches(0.98))
        bar.fill.solid()
        bar.fill.fore_color.rgb = color
        bar.line.fill.background()
        add_text(s, Inches(0.8), y + Inches(0.13), Inches(11.8), Inches(0.32),
                 title, size=14, bold=True, color=color)
        add_text(s, Inches(0.8), y + Inches(0.48), Inches(11.8), Inches(0.45),
                 body, size=12, color=TEXT_MAIN)
        y += Inches(1.06)

    add_footer(s, 17, TOTAL)


# ── Build all ───────────────────────────────────────────────────────────
slide_01_cover()
slide_02_what_is_it()
slide_03_design_dna()
slide_04_prior_art()
slide_05_pipeline()
slide_06_ingest_flow()
slide_07_data_model()
slide_08_query_flow()
slide_09_llm_usage()
slide_10_internal_llm()
slide_11_management()
slide_12_federation()
slide_13_zenithloom()
slide_14_module_map()
slide_15_kg_alignment()
slide_16_roadmap()
slide_17_summary()

out = Path(__file__).parent / "PrismRag-Design-Overview.pptx"
prs.save(out)
print(f"✓ Saved: {out}")
print(f"  Slides: {len(prs.slides)}")
